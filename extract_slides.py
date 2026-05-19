"""
Extrai imagens de fundo do template PPTX para static/slides/.
Executado uma vez durante o docker build.
"""
import zipfile
import os
import hashlib

os.makedirs('static/slides', exist_ok=True)

pptx = 'input_files/cotacao_auto.pptx'
if os.path.exists(pptx):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(pptx)
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                h = hashlib.md5(img.blob).hexdigest()[:8]
                name = f'slide{i+1}_{shape.name.replace(" ","_")}_{h}.{img.ext}'
                path = f'static/slides/{name}'
                if not os.path.exists(path):
                    with open(path, 'wb') as f:
                        f.write(img.blob)
                    print(f'Extracted: {name}')
                break  # apenas a primeira imagem (background) por slide
    print('Extração concluída.')
else:
    print('AVISO: template PPTX não encontrado.')
