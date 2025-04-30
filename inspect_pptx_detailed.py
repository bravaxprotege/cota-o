import sys
sys.path.append('/opt/.manus/.sandbox-runtime')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def inspect_pptx_detailed(pptx_path):
    prs = Presentation(pptx_path)
    print(f"Total de slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n=== SLIDE {i+1} ===")
        print(f"Número de formas: {len(slide.shapes)}")
        
        for j, shape in enumerate(slide.shapes):
            print(f"\nForma {j+1}:")
            print(f"  Nome: {shape.name}")
            print(f"  Tipo: {shape.shape_type}")
            
            # Verificar se é um placeholder
            if shape.is_placeholder:
                ph = shape.placeholder_format
                print(f"  É um placeholder - Índice: {ph.idx}, Tipo: {ph.type}")
            
            # Verificar se tem texto
            if shape.has_text_frame:
                print(f"  Tem frame de texto:")
                text = shape.text_frame.text
                print(f"  Texto: '{text[:100]}{'...' if len(text) > 100 else ''}'")
                
                # Verificar parágrafos
                for k, paragraph in enumerate(shape.text_frame.paragraphs):
                    print(f"    Parágrafo {k+1}: '{paragraph.text[:50]}{'...' if len(paragraph.text) > 50 else ''}'")
            
            # Verificar se é uma tabela
            if shape.has_table:
                table = shape.table
                print(f"  É uma tabela com {len(table.rows)} linhas e {len(table.columns)} colunas")
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        text = cell.text_frame.text
                        print(f"    Célula [{row_idx},{col_idx}]: '{text[:30]}{'...' if len(text) > 30 else ''}'")
            
            # Verificar se é uma imagem
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print("  É uma imagem")

pptx_file = '/home/ubuntu/upload/Cotação auto.pptx'
inspect_pptx_detailed(pptx_file)
