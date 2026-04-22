import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import traceback
import os
import logging
from lxml import etree

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(name)s - %(message)s')
log_prefix = "[preenche_cotacao]"

# Namespace DrawingML
NSMAP = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# ── Formatação exata extraída do template original ───────────────────────────
FONTE_DADOS       = 'Arial Black'        # fonte dos campos de dados (slide 4)
FONTE_PRECO       = 'Arial Black'        # fonte dos valores de preço
TAMANHO_DADOS     = Pt(24)              # tamanho dos campos da ficha do veículo
TAMANHO_PRECO     = Pt(44)              # tamanho padrão dos preços (Ouro/Diamante)
TAMANHO_PRECO_PLAT = Pt(48)            # Platinum tem caixas maiores → 48pt
COR_PRECO = RGBColor(0xFF, 0xC0, 0x00) # #FFC000 — amarelo do template


def format_currency_manual(value):
    """R$ XXX.XXX,XX"""
    if value is None or not isinstance(value, (int, float)):
        return "N/A"
    try:
        return f"R$ {value:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    except (TypeError, ValueError):
        return "Valor inválido"


def format_currency_value_only(value):
    """XXX.XXX,XX  (sem R$)"""
    if value is None or not isinstance(value, (int, float)):
        return "N/A"
    try:
        return f"{value:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    except (TypeError, ValueError):
        return "Valor inválido"


def set_text(text_frame, text_value,
             font_name=FONTE_DADOS,
             font_size=None,
             bold=False,
             color=None,
             alignment=None,
             word_wrap=False):
    """
    Escreve texto com formatação 100% controlada.
    - word_wrap=False: evita quebra de linha (padrão para todos os campos).
    - Reseta espaçamento de caracteres e kerning via XML.
    - Zera espaçamento de parágrafo (antes/depois).
    - Cor, fonte e tamanho aplicados explicitamente.
    """
    if text_frame is None:
        return

    # Lê alinhamento do template antes de apagar
    tmpl_alignment = alignment
    if tmpl_alignment is None and text_frame.paragraphs:
        tmpl_alignment = text_frame.paragraphs[0].alignment

    tamanho = font_size if font_size is not None else TAMANHO_DADOS

    # Desativa quebra de linha
    try:
        text_frame.word_wrap = word_wrap
    except Exception:
        pass

    text_frame.clear()
    p = text_frame.add_paragraph()

    if tmpl_alignment is not None:
        try:
            p.alignment = tmpl_alignment
        except Exception as e:
            logging.warning(f"  Alinhamento não aplicado: {e}")

    # Zera espaçamento do parágrafo via XML (spcBef / spcAft)
    try:
        pPr = p._p.get_or_add_pPr()
        # Remove espaçamento entre linhas herdado do template
        for tag in (f'{{{NSMAP}}}lnSpc',):
            el = pPr.find(tag)
            if el is not None:
                pPr.remove(el)
        pPr.set('spcBef', '0')
        pPr.set('spcAft', '0')
    except Exception as e:
        logging.warning(f"  Espaçamento de parágrafo não resetado: {e}")

    run = p.add_run()
    run.text = str(text_value)
    run.font.name  = font_name
    run.font.size  = tamanho
    run.font.bold  = bold

    if color is not None:
        run.font.color.rgb = color

    # Reset COMPLETO de espaçamento de caracteres via XML
    try:
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', '0')               # Espaço extra entre caracteres = 0
        rPr.attrib.pop('kern', None)       # Remove kerning (quando ausente = desativado)
        rPr.attrib.pop('baseline', None)   # Remove deslocamento de linha base
    except Exception as e:
        logging.warning(f"  Espaçamento de run não resetado: {e}")

    logging.info(f"  '{text_value}' → {font_name} {tamanho.pt}pt bold={bold} cor={color}")


def remover_slide(prs, slide_index):
    """Remove slide pelo índice. Sempre do maior para o menor."""
    xml_slides = prs.slides._sldIdLst
    rId = xml_slides[slide_index].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    prs.part.drop_rel(rId)
    del xml_slides[slide_index]
    logging.info(f"{log_prefix} Slide {slide_index} removido.")


def preencher_cotacao_pptx(template_path, output_path, dados_cotacao):
    logging.info(f"{log_prefix} Iniciando: {template_path}")
    try:
        prs = Presentation(template_path)
        logging.info(f"{log_prefix} Template aberto. Slides: {len(prs.slides)}")

        def find_shape(slide_index, shape_name):
            if slide_index < 0 or slide_index >= len(prs.slides):
                logging.warning(f"{log_prefix} Slide {slide_index} não existe.")
                return None
            for shape in prs.slides[slide_index].shapes:
                if shape.name.strip().lower() == shape_name.strip().lower():
                    if shape.has_text_frame:
                        return shape.text_frame
                    logging.warning(f"{log_prefix} '{shape.name}' sem text frame.")
                    return None
            logging.warning(f"{log_prefix} Shape '{shape_name}' não encontrada no slide {slide_index+1}.")
            return None

        # ── Dados ───────────────────────────────────────────────────────────
        nome_cliente      = dados_cotacao.get("nome_cliente", "N/A")
        placa             = dados_cotacao.get("placa", "N/A").upper()
        marca             = dados_cotacao.get("marca", "N/A").upper()
        modelo            = dados_cotacao.get("modelo", "N/A").upper()
        ano               = dados_cotacao.get("ano", "N/A")
        valor_fipe        = dados_cotacao.get("valor_fipe")
        categoria         = dados_cotacao.get("categoria", "")
        precos            = dados_cotacao.get("precos", {})
        veiculo_pesado    = dados_cotacao.get("veiculo_pesado", False)
        sujeito_aprovacao = precos.get("sujeito_aprovacao", False)
        adesao_str        = format_currency_value_only(precos.get('Adesão'))

        # ── Slide 1 / índice 0 — Capa ───────────────────────────────────────
        logging.info(f"{log_prefix} Slide 1 (capa)")
        set_text(find_shape(0, "Nome associado"),
                 nome_cliente.title(),
                 font_size=Pt(20), bold=True)

        # ── Slide 4 / índice 3 — Ficha do veículo ───────────────────────────
        logging.info(f"{log_prefix} Slide 4 (ficha)")
        set_text(find_shape(3, "Nome associado"), nome_cliente.title(), bold=True)
        set_text(find_shape(3, "Placa"),          placa,                bold=True)
        set_text(find_shape(3, "Marca carro"),    marca,                bold=True)
        set_text(find_shape(3, "modelo"),         modelo,               bold=True)
        set_text(find_shape(3, "Ano"),            str(ano),             bold=True)
        set_text(find_shape(3, "Categoria"),
                 categoria.upper() if categoria else "PASSEIO",         bold=True)
        set_text(find_shape(3, "Valor fipe"),
                 format_currency_manual(valor_fipe),                    bold=True)

        if veiculo_pesado:
            # ── PESADO: apenas Diamante (índice 5) com cor amarela ──────────
            logging.info(f"{log_prefix} Modo PESADO")
            preco_pesado_str = format_currency_value_only(precos.get('Pesados'))

            set_text(find_shape(5, "adesão"),
                     adesao_str, bold=True,
                     font_size=TAMANHO_PRECO, color=COR_PRECO)
            set_text(find_shape(5, "diamante"),
                     preco_pesado_str, bold=True,
                     font_size=TAMANHO_PRECO, color=COR_PRECO)

            # Remove Platinum (6) e Ouro (4) — maior primeiro
            total = len(prs.slides)
            if total > 6: remover_slide(prs, 6)
            if total > 4: remover_slide(prs, 4)

        else:
            # ── NORMAL: 3 planos com cor amarela ────────────────────────────
            logging.info(f"{log_prefix} Modo NORMAL")

            # Ouro — índice 4 — 44pt
            set_text(find_shape(4, "adesão"),
                     adesao_str, bold=True,
                     font_size=TAMANHO_PRECO, color=COR_PRECO)
            set_text(find_shape(4, "ouro"),
                     format_currency_value_only(precos.get('Plano Ouro')),
                     bold=True, font_size=TAMANHO_PRECO, color=COR_PRECO)

            # Diamante — índice 5 — 44pt
            set_text(find_shape(5, "adesão"),
                     adesao_str, bold=True,
                     font_size=TAMANHO_PRECO, color=COR_PRECO)
            set_text(find_shape(5, "diamante"),
                     format_currency_value_only(precos.get('Diamante')),
                     bold=True, font_size=TAMANHO_PRECO, color=COR_PRECO)

            # Platinum — índice 6 — 48pt (caixas maiores no template)
            set_text(find_shape(6, "adesão"),
                     adesao_str, bold=True,
                     font_size=TAMANHO_PRECO_PLAT, color=COR_PRECO)
            set_text(find_shape(6, "platinium"),
                     format_currency_value_only(precos.get('Platinum')),
                     bold=True, font_size=TAMANHO_PRECO_PLAT, color=COR_PRECO)

        if sujeito_aprovacao:
            logging.warning(f"{log_prefix} Sujeito à aprovação — shape não definida.")

        prs.save(output_path)
        logging.info(f"{log_prefix} Salvo: {output_path}")
        return True

    except Exception as e:
        logging.error(f"{log_prefix} ERRO: {e}")
        traceback.print_exc()
        return False


if __name__ == "__main__":
    pass
