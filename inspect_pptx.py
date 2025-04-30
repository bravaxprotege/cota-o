import sys
sys.path.append('/opt/.manus/.sandbox-runtime')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def find_placeholders(pptx_path):
    prs = Presentation(pptx_path)
    placeholders = {}
    for i, slide in enumerate(prs.slides):
        slide_placeholders = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph = shape.placeholder_format
                slide_placeholders.append(f"Slide {i+1}, Placeholder Index: {ph.idx}, Type: {ph.type}, Shape Name: {shape.name}")
            elif shape.has_text_frame:
                # Check for potential text boxes used as placeholders (e.g., {{FieldName}})
                if '{{' in shape.text_frame.text and '}}' in shape.text_frame.text:
                     slide_placeholders.append(f"Slide {i+1}, Text Box (Potential Placeholder): {shape.name}, Text: '{shape.text_frame.text[:50]}...' ")
            # You might need to inspect other shape types if they contain dynamic data

        if slide_placeholders:
            placeholders[f'Slide_{i+1}'] = slide_placeholders
    return placeholders

pptx_file = '/home/ubuntu/upload/Cotação auto.pptx'
found_placeholders = find_placeholders(pptx_file)

# Print the found placeholders
for slide, ph_list in found_placeholders.items():
    print(f"--- {slide} ---")
    for ph in ph_list:
        print(ph)

# Also, let's try to list all text frames to catch non-placeholder text that might need changing
print("\n--- All Text Frames ---")
prs = Presentation(pptx_file)
for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(f"Shape Name: {shape.name}, Text: '{shape.text_frame.text[:100]}...' ")

